from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_poster():
    prs = Presentation()
    # Définir la taille de la diapositive (Format A4 Vertical pour un poster standard ou A3)
    # Ici on met un format A3 vertical pour avoir de la place (29.7cm x 42cm)
    prs.slide_width = Cm(29.7)
    prs.slide_height = Cm(42)

    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 6 = Blank layout

    # --- COULEURS ---
    BLUE_DARK = RGBColor(31, 73, 125)
    ORANGE_LIGHT = RGBColor(252, 228, 214) # Fond beige/orange clair
    ORANGE_TITLE = RGBColor(237, 125, 49)  # Texte orange
    GREY_BG = RGBColor(242, 242, 242)

    # --- FONCTION D'AIDE POUR CRÉER DES BLOCS ---
    def add_section_box(slide, left, top, width, height, title_text, bg_color=None):
        # Fond de section
        if bg_color:
            shape = slide.shapes.add_shape(1, left, top, width, height) # 1 = Rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.color.rgb = RGBColor(200, 200, 200)
        
        # Titre de la section
        textbox = slide.shapes.add_textbox(left, top, width, Cm(1.5))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = ORANGE_TITLE

        # Zone de texte modifiable en dessous
        content_box = slide.shapes.add_textbox(left, top + Cm(1), width, height - Cm(1))
        tf_content = content_box.text_frame
        tf_content.text = "Ajouter votre texte ici..."
        return content_box

    # --- 1. EN-TÊTE (HEADER) ---
    header_h = Cm(4)
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_DARK
    header.line.fill.background()

    # Titre Principal
    title_box = slide.shapes.add_textbox(0, Cm(0.5), prs.slide_width, Cm(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DRONEA"
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    sub_box = slide.shapes.add_textbox(0, Cm(3), prs.slide_width, Cm(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Interface de téléopération intuitive pour robot mobile avec retour vidéo FPV"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Info Étudiants (Barre grise)
    info_bar = slide.shapes.add_shape(1, 0, Cm(4.2), prs.slide_width, Cm(2))
    info_bar.fill.solid()
    info_bar.fill.fore_color.rgb = GREY_BG
    info_bar.line.fill.background()
    
    info_text = slide.shapes.add_textbox(Cm(1), Cm(4.5), prs.slide_width - Cm(2), Cm(1.5))
    info_text.text_frame.text = "Kevin Tran, Thomas X, Romain X... | EFREI Paris\nProjet entreprise - Sigmatics"

    # --- COLONNE GAUCHE ---
    col_gap = Cm(0.5)
    col_width = (prs.slide_width - (Cm(1)*2) - col_gap) / 2
    left_x = Cm(1)
    right_x = left_x + col_width + col_gap
    current_y = Cm(6.5)

    # Section 1: Fil conducteur
    add_section_box(slide, left_x, current_y, col_width, Cm(5), "1. Fil conducteur", ORANGE_LIGHT)
    current_y += Cm(5.5)

    # Section 3: Solution proposée (partie gauche)
    add_section_box(slide, left_x, current_y, col_width, Cm(4), "3. Solution proposée")
    current_y += Cm(4.5)

    # Section 4: Objectifs du projet
    add_section_box(slide, left_x, current_y, col_width, Cm(3), "4. Objectifs du projet", ORANGE_LIGHT)
    current_y += Cm(3.5)

    # Section 5: Logique de pilotage
    add_section_box(slide, left_x, current_y, col_width, Cm(4), "5. Logique de pilotage")
    current_y += Cm(4.5)
    
    # Section 8: Composants techniques
    add_section_box(slide, left_x, current_y, col_width, Cm(4), "8. Composants techniques", GREY_BG)
    
    # --- COLONNE DROITE ---
    current_y = Cm(6.5) # Reset Y

    # Section 2: Solution proposée (partie droite)
    add_section_box(slide, right_x, current_y, col_width, Cm(4.5), "2. Solution proposée", ORANGE_LIGHT)
    current_y += Cm(5)

    # Section 4: Architecture & Flux (Le schéma compliqué)
    add_section_box(slide, right_x, current_y, col_width, Cm(6), "4. Architecture & flux de données")
    # Placeholder pour l'image du schéma
    pic_ph = slide.shapes.add_shape(1, right_x + Cm(1), current_y + Cm(1.5), col_width - Cm(2), Cm(4))
    pic_ph.text_frame.text = "[Insérer Schéma Manette -> Raspberry ici]"
    current_y += Cm(6.5)

    # Section 7: Retour vidéo & HUD
    add_section_box(slide, right_x, current_y, col_width, Cm(6), "7. Retour vidéo & interface HUD")
    # Placeholder pour l'image
    pic_ph2 = slide.shapes.add_shape(1, right_x + Cm(1), current_y + Cm(1.5), col_width - Cm(2), Cm(4))
    pic_ph2.text_frame.text = "[Insérer Interface Graphique ici]"
    current_y += Cm(6.5)

    # Section 11: Conclusion
    add_section_box(slide, right_x, current_y, col_width, Cm(4), "11. Conclusion & perspectives")

    # --- PIED DE PAGE (TOTAL) ---
    footer_y = prs.slide_height - Cm(2)
    footer = slide.shapes.add_shape(1, 0, footer_y, prs.slide_width, Cm(2))
    footer.fill.solid()
    footer.fill.fore_color.rgb = BLUE_DARK
    
    footer_text = slide.shapes.add_textbox(Cm(1), footer_y + Cm(0.5), prs.slide_width, Cm(1))
    p = footer_text.text_frame.paragraphs[0]
    p.text = "TOTAL    ≈ 50 €"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ORANGE_TITLE

    # Sauvegarde
    prs.save('poster_dronea_recreated.pptx')
    print("Poster généré avec succès : poster_dronea_recreated.pptx")

create_poster()